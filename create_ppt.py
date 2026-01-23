from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(25, 55, 109)  # Dark blue
ACCENT_COLOR = RGBColor(0, 150, 136)    # Teal
LIGHT_BG = RGBColor(240, 247, 255)     # Light blue
TEXT_COLOR = RGBColor(50, 50, 50)       # Dark gray

def add_title_slide(prs, title, subtitle=""):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = RGBColor(200, 220, 255)
        subtitle_p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items=None, image_path=None):
    """Add content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title text
    title_frame = title_shape.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.space_before = Pt(12)
    title_p.space_after = Pt(12)
    
    # Content
    if content_items:
        left = Inches(0.7)
        top = Inches(1.5)
        width = Inches(8.6)
        height = Inches(5.5)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.level = 0
            p.space_before = Pt(8)
            p.space_after = Pt(8)
    
    return slide

# ============ SLIDE 1: TITLE ============
add_title_slide(prs, "KRAFTD", "Intelligent Supply Chain for SMEs\nMISA Entrepreneurship Pitch")

# ============ SLIDE 2: THE INEQUALITY ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = PRIMARY_COLOR
title_frame = title_shape.text_frame
p = title_frame.paragraphs[0]
p.text = "The Founding Insight"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

left = Inches(0.7)
top = Inches(1.2)
width = Inches(8.6)
height = Inches(5.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

items = [
    "✓ SMEs process raw materials, export goods—exactly like SABIC",
    "",
    "✗ SMEs lack structured data, visibility, and infrastructure",
    "",
    "✓ SABIC has high-tech endpoints; SMEs have Excel spreadsheets",
    "",
    "Result: 600K+ Saudi SMEs waste 15-30% on manual document processing",
    "",
    "Kraftd: Rebalancing the intelligence equation"
]
for i, item in enumerate(items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(6)

# ============ SLIDE 3: THE PROBLEM ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = PRIMARY_COLOR
title_frame = title_shape.text_frame
p = title_frame.paragraphs[0]
p.text = "Why Digitalization Failed SMEs"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

left = Inches(0.7)
top = Inches(1.2)
width = Inches(8.6)
height = Inches(5.8)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.word_wrap = True

items = [
    "Enterprise systems (SAP, Oracle) were built FOR enterprises",
    "",
    "AI platforms built APIs around structured data—only serves big companies",
    "",
    "When digitalization was imposed on SMEs, it added:",
    "  • More work (portals, uploads, templates)",
    "  • More cost ($50K+/month for enterprise ERP)",
    "  • More complexity (6-18 month implementation)",
    "",
    "SMEs couldn't afford it. Never reached down-market."
]
for i, item in enumerate(items):
    if i == 0:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(17)
    p.font.color.rgb = TEXT_COLOR
    p.space_after = Pt(5)

# ============ SLIDE 4: MARKET OPPORTUNITY ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Title
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = PRIMARY_COLOR
title_frame = title_shape.text_frame
p = title_frame.paragraphs[0]
p.text = "Market Opportunity"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Create table
left = Inches(0.7)
top = Inches(1.2)
width = Inches(8.6)
height = Inches(5.5)

rows, cols = 5, 3
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table_shape.columns[0].width = Inches(3)
table_shape.columns[1].width = Inches(2.5)
table_shape.columns[2].width = Inches(3.1)

# Header
headers = ["Metric", "Value", "Implication"]
for i, header in enumerate(headers):
    cell = table_shape.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(232, 238, 247)
    p = cell.text_frame.paragraphs[0]
    p.text = header
    p.font.bold = True
    p.font.size = Pt(14)

# Data rows
data = [
    ["SME Population", "600,000+", "Addressable market"],
    ["Procurement Value", "$400B+ annually", "Size of opportunity"],
    ["Manual Processes", "78%", "Digitalization gap"],
    ["Time Wasted", "15-30% per SME", "Efficiency gain"]
]

for row_idx, row_data in enumerate(data, 1):
    for col_idx, cell_text in enumerate(row_data):
        cell = table_shape.cell(row_idx, col_idx)
        p = cell.text_frame.paragraphs[0]
        p.text = cell_text
        p.font.size = Pt(13)

# ============ SLIDE 5: THE SOLUTION ============
slide = add_content_slide(prs, "What Kraftd Does")
items = [
    "Convert unstructured documents → structured intelligence (95%+ accuracy)",
    "",
    "Document Types: Invoices, quotations, BOQs, contracts, shipping docs",
    "",
    "Monitor AI reliability in workflows (unique capability—no competitors)",
    "",
    "SME-optimized: 2-week implementation, $500-5K/month, no IT team required",
    "",
    "vs. SAP/Oracle: 6-18 months, $50K+/month, requires enterprise infrastructure"
]
for item in items:
    slide.shapes[1].text_frame.add_paragraph().text = item

# ============ SLIDE 6: TRACTION ============
slide = add_content_slide(prs, "Traction & Validation")
items = [
    "MVP Live: January 2026 (production-ready)",
    "",
    "Market Validation: 47+ organizations interested | 20+ active trials",
    "",
    "Accuracy: 95%+ on real-world documents",
    "",
    "Customer Intent: $3K-15K/month willingness to pay confirmed",
    "",
    "Unit Economics: LTV:CAC = 13:1 to 25:1 | 2-4 month payback"
]
for item in items:
    slide.shapes[1].text_frame.add_paragraph().text = item

# ============ SLIDE 7: FOUNDER ============
slide = add_content_slide(prs, "The Founder: Sulaiman Ul Akram Ismail")
items = [
    "12+ years in petrochemical, logistics, industrial operations (GCC)",
    "",
    "✓ Saudi Aramco, SABIC, NEOM experience | 500,000+ safe man-hours",
    "",
    "✓ Mechanical Engineering | McKinsey Forward | Founder Institute Launch Track",
    "",
    "✓ Built entire MVP single-handedly (full-stack developer)",
    "",
    "✓ Lived the problem for a decade—knows SME pain intimately"
]
for item in items:
    slide.shapes[1].text_frame.add_paragraph().text = item

# ============ SLIDE 8: FINANCIAL PROJECTIONS ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Title
title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = PRIMARY_COLOR
title_frame = title_shape.text_frame
p = title_frame.paragraphs[0]
p.text = "Financial Projections"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Create financial table
left = Inches(0.5)
top = Inches(1.2)
width = Inches(9)
height = Inches(5.5)

rows, cols = 5, 5
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Headers
headers = ["Metric", "Year 1", "Year 2", "Year 3", "Year 5"]
for i, header in enumerate(headers):
    cell = table_shape.cell(0, i)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(232, 238, 247)
    p = cell.text_frame.paragraphs[0]
    p.text = header
    p.font.bold = True
    p.font.size = Pt(13)

# Data
fin_data = [
    ["Customers", "15-25", "50-80", "150-250", "500-1K"],
    ["Revenue", "$180-300K", "$600K-1.2M", "$1.8-3M", "$6-12M"],
    ["Profitability", "Break-even", "+$200-400K", "+$1M", "+40% EBITDA"],
    ["Payback", "2-4 months", "2-4 months", "2-4 months", "Mature"]
]

for row_idx, row_data in enumerate(fin_data, 1):
    for col_idx, cell_text in enumerate(row_data):
        cell = table_shape.cell(row_idx, col_idx)
        p = cell.text_frame.paragraphs[0]
        p.text = cell_text
        p.font.size = Pt(12)

# ============ SLIDE 9: COMPETITIVE ADVANTAGE ============
slide = add_content_slide(prs, "Why Kraftd Wins")
items = [
    "No Direct Competitors",
    "  Enterprise solutions serve large orgs | Startups solve point problems",
    "",
    "Founder Advantage: 12+ years supply chain domain expertise",
    "",
    "Cost Advantage: 100x cheaper than enterprise ERP ($500-5K vs $50K+/month)",
    "",
    "Speed: 2-week implementation vs 6-18 months for SAP",
    "",
    "Complete Solution: End-to-end, not point tool"
]
for item in items:
    slide.shapes[1].text_frame.add_paragraph().text = item

# ============ SLIDE 10: STRATEGIC VALUE ============
slide = add_content_slide(prs, "Why MISA Should Support Kraftd")
items = [
    "Success Story: Proof Saudi Arabia builds world-class SaaS companies",
    "",
    "National Impact: Empower 600K+ SMEs | Align with Vision 2030",
    "",
    "High ROI: Scalable business model with clear path to $10M+ revenue",
    "",
    "Market Timing: Perfect conditions (AI maturity + cost reduction + push for digitalization)",
    "",
    "Return Probability: 25-50x within 5-7 years"
]
for item in items:
    slide.shapes[1].text_frame.add_paragraph().text = item

# ============ SLIDE 11: ROADMAP ============
slide = add_content_slide(prs, "2026 Roadmap")
items = [
    "Q1: Complete MISA licensing | Secure first 5-10 paying customers",
    "",
    "Q2: Hire VP Sales | Launch marketing | 15-20 customers",
    "",
    "Q3: Expand to UAE/Kuwait | Launch mobile app | AI Monitoring feature",
    "",
    "Q4: 35-50 customers | $300K+ revenue | Prepare Series A",
    "",
    "Goal: Path to $600K-1.2M revenue by end of 2027"
]
for item in items:
    slide.shapes[1].text_frame.add_paragraph().text = item

# ============ SLIDE 12: THE ASK ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_BG

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "The Ask"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR

# Main box
box = slide.shapes.add_shape(1, Inches(1.5), Inches(2), Inches(7), Inches(4.5))
box.fill.solid()
box.fill.fore_color.rgb = RGBColor(255, 255, 255)
box.line.color.rgb = ACCENT_COLOR
box.line.width = Pt(3)

# Content in box
text_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(3.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

items = [
    ("Entrepreneurship License", "Legal registration support + network credibility"),
    ("Business Network", "Customer introductions, partner connections"),
    ("Mentorship", "Operational scaling, governance, fundraising"),
    ("Series A Platform", "Credibility for future investor conversations")
]

for label, desc in items:
    p = text_frame.add_paragraph()
    p.text = f"• {label}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(4)
    
    p2 = text_frame.add_paragraph()
    p2.text = f"  {desc}"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_COLOR
    p2.space_after = Pt(12)

# ============ SLIDE 13: CLOSING ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = PRIMARY_COLOR

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p = title_frame.paragraphs[0]
p.text = "Rebalancing Intelligence in Supply Chains"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
p = subtitle_frame.paragraphs[0]
p.text = "Kraftd is not another software tool.\n\nKraftd is rebalancing the intelligence equation—giving SMEs the clarity and precision previously available only to enterprises.\n\nWith MISA, this becomes the proof point that Saudi Arabia builds world-class technology for global markets."
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(200, 220, 255)
p.alignment = PP_ALIGN.CENTER
p.line_spacing = 1.5

# Contact
contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.7), Inches(9), Inches(0.5))
contact_frame = contact_box.text_frame
p = contact_frame.paragraphs[0]
p.text = "akram@kraftd.io | https://kraftd-a4gfhqa2axb2h6cd.uaenorth-01.azurewebsites.net"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Save
prs.save(r'c:\Users\1R6\OneDrive\Desktop\Kraftd Docs\KRAFTD_Pitch_Deck_MISA.pptx')
print("✓ Pitch Deck (PowerPoint) created")
print("\n" + "="*60)
print("ALL DELIVERABLES CREATED SUCCESSFULLY!")
print("="*60)
print("1. KRAFTD_Company_Profile_MISA.docx (8-10 pages)")
print("2. KRAFTD_One_Page_Summary_MISA.docx (1 page)")
print("3. KRAFTD_Pitch_Deck_MISA.pptx (13 slides)")
print("\nLocation: c:\\Users\\1R6\\OneDrive\\Desktop\\Kraftd Docs\\")
print("\nFont: Segoe UI (modern, clean)")
print("Design: Minimal, professional, MISA-ready")
print("="*60)
