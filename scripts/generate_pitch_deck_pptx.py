#!/usr/bin/env python3
"""
KRAFTD Strategic Pitch Deck - PowerPoint Generator
Produces a professional, investor-grade presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# Brand Colors
COLOR_PRIMARY = RGBColor(0, 51, 102)      # Dark blue
COLOR_ACCENT = RGBColor(255, 102, 0)      # Orange
COLOR_SECONDARY = RGBColor(51, 102, 153)  # Light blue
COLOR_TEXT = RGBColor(51, 51, 51)         # Dark gray
COLOR_WHITE = RGBColor(255, 255, 255)     # White
COLOR_LIGHT_GRAY = RGBColor(240, 240, 240) # Light gray

def create_presentation():
    """Create the main presentation object"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs, title, subtitle):
    """Add a professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_box.text_frame
    p = date_frame.paragraphs[0]
    p.text = "January 2026 | Series A Investment"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        
        p = text_frame.paragraphs[i]
        
        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 20))
            p.font.bold = item.get('bold', False)
            p.font.color.rgb = item.get('color', COLOR_TEXT)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXT
        
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column comparison slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Left column
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.4))
    left_frame = left_header.text_frame
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    left_text = left_box.text_frame
    left_text.word_wrap = True
    
    for i, item in enumerate(left_items):
        if i > 0:
            left_text.add_paragraph()
        p = left_text.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.4))
    right_frame = right_header.text_frame
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    right_text = right_box.text_frame
    right_text.word_wrap = True
    
    for i, item in enumerate(right_items):
        if i > 0:
            right_text.add_paragraph()
        p = right_text.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

def add_chart_slide(prs, title, chart_title, data_labels, data_values, chart_type="bar"):
    """Add a slide with a chart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = data_labels
    
    if chart_type == "bar":
        chart_data.add_series('Value', (data_values,))
        x, y, cx, cy = Inches(1), Inches(1.2), Inches(8), Inches(5.8)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
    
    chart.has_legend = False
    chart.chart_title.text_frame.text = ""
    
    return slide

def add_financial_table_slide(prs):
    """Add financial projection slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Top bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.color.rgb = COLOR_PRIMARY
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "5-Year Financial Projection"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Table
    rows, cols = 8, 6
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.8)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ['Metric', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5']
    for i, header in enumerate(headers):
        cell = table_shape.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = COLOR_WHITE
    
    # Data rows
    data = [
        ['Customers', '560', '1,200', '2,500', '4,000', '5,500+'],
        ['ARR', '$5.8M', '$18.2M', '$58M', '$92M', '$135M+'],
        ['Revenue', '$6.2M', '$19.8M', '$62M', '$100M', '$147M+'],
        ['EBITDA', '$3.4M', '$12.2M', '$42.5M', '$68.5M', '$100.3M'],
        ['EBITDA %', '55%', '62%', '69%', '69%', '68%'],
        ['FCF', '$2.8M', '$10.5M', '$38M', '$62M', '$90M+'],
        ['Status', '✓ Profitable', '↑ Scaling', '↑ Growth', '↑ Expansion', '↑ Scale'],
    ]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_value in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_value
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    if col_idx == 0:
                        run.font.bold = True
    
    return slide

# Create presentation
prs = create_presentation()

print("Generating KRAFTD Strategic Pitch Deck...")

# SLIDE 1: Title Slide
print("  → Slide 1: Title Slide")
add_title_slide(prs, "KRAFTD", "Intelligence Equity in Supply Chains")

# SLIDE 2: The Problem
print("  → Slide 2: The Problem")
add_content_slide(prs, "The Problem", [
    {'text': 'Enterprise Supply Chains', 'size': 24, 'bold': True, 'color': COLOR_PRIMARY},
    {'text': 'ERP systems, clean data, APIs, dashboards, AI/ML', 'size': 18, 'level': 1},
    {'text': 'SME Supply Chains', 'size': 24, 'bold': True, 'color': COLOR_ACCENT},
    {'text': 'Spreadsheets, emails, manual tracking, zero intelligence', 'size': 18, 'level': 1},
    {'text': '', 'size': 14},
    {'text': '80% of supply chain operators have zero intelligent visibility', 'size': 20, 'bold': True, 'color': COLOR_PRIMARY},
])

# SLIDE 3: Market Opportunity
print("  → Slide 3: Market Opportunity")
add_content_slide(prs, "The Opportunity", [
    {'text': 'Global Supply Chain: $118 Trillion/Year', 'size': 22, 'bold': True},
    {'text': 'Enterprise-served (Coupa, SAP): $500B', 'size': 18, 'level': 1, 'color': COLOR_SECONDARY},
    {'text': 'SME/Supplier segment (UNSERVED): $2.5 Trillion', 'size': 18, 'level': 1, 'color': COLOR_ACCENT},
    {'text': '', 'size': 14},
    {'text': 'GCC Region Addressable Market:', 'size': 22, 'bold': True},
    {'text': 'Saudi Arabia: $550B | UAE: $300B | Others: $600B', 'size': 18, 'level': 1},
    {'text': 'Total GCC: $1.45 Trillion', 'size': 20, 'bold': True, 'color': COLOR_PRIMARY},
])

# SLIDE 4: Why Incumbents Fail
print("  → Slide 4: Why Incumbents Cannot Compete")
add_two_column_slide(prs, "Why Incumbents Cannot Compete",
    "What Coupa Requires", [
        "✗ ERP System",
        "✗ Clean data",
        "✗ IT integration team",
        "✗ $500K-$1M budget",
        "✗ 3-6 month deployment",
        "✗ Min. $50M revenue",
        "→ Only 0.3% of supply chain operators"
    ],
    "What KRAFTD Requires", [
        "✓ Browser access",
        "✓ Any documents",
        "✓ $5K-$50K/month",
        "✓ 2-4 week onboarding",
        "✓ Any company size",
        "✓ No prerequisites",
        "→ 85% of supply chain operators"
    ]
)

# SLIDE 5: The Solution
print("  → Slide 5: The Solution")
add_content_slide(prs, "The Solution: Translation Layer", [
    {'text': 'Takes messy reality...', 'size': 20, 'bold': True},
    {'text': 'PDFs, invoices, contracts, customs forms, emails', 'size': 16, 'level': 1},
    {'text': '↓', 'size': 20, 'bold': True},
    {'text': 'Extracts structured intelligence...', 'size': 20, 'bold': True},
    {'text': '95%+ accuracy, actionable insights', 'size': 16, 'level': 1},
    {'text': '↓', 'size': 20, 'bold': True},
    {'text': 'Provides decision support...', 'size': 20, 'bold': True},
    {'text': 'Price trends, supplier reliability, risk indicators, competitive positioning', 'size': 16, 'level': 1},
])

# SLIDE 6: Product-Market Fit
print("  → Slide 6: Product-Market Fit")
add_content_slide(prs, "Product-Market Fit Proven", [
    {'text': 'Document Accuracy: 95.3% (industry benchmark: 88-92%)', 'size': 18},
    {'text': 'Implementation: 2-3 weeks (vs Coupa: 4-6 months)', 'size': 18},
    {'text': 'Cost: $5-50K (vs Coupa: $500K-$1.5M)', 'size': 18},
    {'text': '', 'size': 12},
    {'text': 'Customer Validation:', 'size': 20, 'bold': True},
    {'text': '47+ expressions of interest in beta', 'size': 16, 'level': 1},
    {'text': 'Early adopters showing ROI within 8 weeks', 'size': 16, 'level': 1},
    {'text': 'Pricing tiers validated', 'size': 16, 'level': 1},
])

# SLIDE 7: Market Validation
print("  → Slide 7: Market Validation")
add_content_slide(prs, "Market Validation", [
    {'text': 'Year 1 Revenue Target: $6.2M', 'size': 22, 'bold': True, 'color': COLOR_ACCENT},
    {'text': 'Procurement-focused SMEs: 625 customers @ $499-2,499/month', 'size': 16, 'level': 1},
    {'text': 'Trade & Import/Export: 64 customers @ $2,499-9,999/month', 'size': 16, 'level': 1},
    {'text': 'Manufacturing/Processing: 72 customers @ $1,499-5,999/month', 'size': 16, 'level': 1},
    {'text': '', 'size': 14},
    {'text': 'Long-term Market:', 'size': 20, 'bold': True},
    {'text': 'GCC Serviceable Addressable Market: $1.26 Trillion', 'size': 16, 'level': 1},
    {'text': 'Year 5-7 Revenue Potential: $6.3B-$12.6B', 'size': 16, 'level': 1},
])

# SLIDE 8: Unit Economics
print("  → Slide 8: Unit Economics")
add_content_slide(prs, "Unit Economics (Why This Works)", [
    {'text': 'Blended LTV:CAC Ratio: 5.3x', 'size': 24, 'bold': True, 'color': COLOR_ACCENT},
    {'text': 'Benchmark excellent: >3x is great', 'size': 14, 'level': 1},
    {'text': '', 'size': 12},
    {'text': 'Blended CAC: $2,850', 'size': 20, 'bold': True},
    {'text': 'Blended LTV (24 months): $15,240', 'size': 20, 'bold': True},
    {'text': 'Payback Period: 3.7 months', 'size': 20, 'bold': True},
    {'text': 'Gross Margin: 91%', 'size': 20, 'bold': True},
    {'text': '', 'size': 12},
    {'text': 'Implication: Every $1 CAC generates $5.30 LTV', 'size': 16, 'italic': True, 'color': COLOR_PRIMARY},
])

# SLIDE 9: Financial Projections
print("  → Slide 9: Financial Projections")
add_financial_table_slide(prs)

# SLIDE 10: GTM Strategy
print("  → Slide 10: Go-to-Market Strategy")
add_content_slide(prs, "Go-to-Market Strategy (Year 1)", [
    {'text': '45% Direct Sales', 'size': 20, 'bold': True},
    {'text': '252 customers | $2,500 CAC | Sales team: 2 AEs + 1 SDR', 'size': 14, 'level': 1},
    {'text': '30% Freemium/Trial', 'size': 20, 'bold': True},
    {'text': '168 customers | $800 CAC | Self-serve path', 'size': 14, 'level': 1},
    {'text': '15% Partnerships', 'size': 20, 'bold': True},
    {'text': '84 customers | $1,200 CAC | System integrator channels', 'size': 14, 'level': 1},
    {'text': '10% Content/Marketing', 'size': 20, 'bold': True},
    {'text': '56 customers | $600 CAC | Inbound demand', 'size': 14, 'level': 1},
])

# SLIDE 11: Competitive Positioning
print("  → Slide 11: Competitive Positioning")
add_two_column_slide(prs, "Competitive Positioning",
    "COUPA POSITION", [
        "Market: Enterprise",
        "Size: $500B",
        "Customers: Fortune 500+",
        "Competition: High",
        "Price: $500K+ implementation"
    ],
    "KRAFTD POSITION", [
        "Market: SME/Supplier",
        "Size: $2.5 Trillion",
        "Customers: All companies",
        "Competition: Zero",
        "Price: $5-50K/month"
    ]
)

# SLIDE 12: The Manifesto
print("  → Slide 12: The Manifesto")
add_content_slide(prs, "The Manifesto", [
    {'text': 'Intelligence should not require ERP maturity.', 'size': 20, 'bold': True},
    {'text': 'Insight should not require API readiness.', 'size': 20, 'bold': True},
    {'text': 'Credibility should not require scale.', 'size': 20, 'bold': True},
    {'text': '', 'size': 14},
    {'text': 'If an SME can operate a business, it should be able to operate with intelligence.', 'size': 20, 'bold': True, 'color': COLOR_ACCENT},
])

# SLIDE 13: Investment Terms
print("  → Slide 13: Investment Terms")
add_content_slide(prs, "Investment Opportunity", [
    {'text': 'Series A: $2-3M', 'size': 28, 'bold': True, 'color': COLOR_ACCENT},
    {'text': '', 'size': 12},
    {'text': 'Use of Funds:', 'size': 20, 'bold': True},
    {'text': 'Sales & GTM: $850K | Marketing: $450K | Product: $350K', 'size': 14, 'level': 1},
    {'text': 'Customer Success: $300K | Operations: $200K | Buffer: $200K', 'size': 14, 'level': 1},
    {'text': '', 'size': 12},
    {'text': 'Return Scenarios:', 'size': 20, 'bold': True},
    {'text': 'Base Case (Year 5): 25-35x return | $2-3B exit valuation', 'size': 14, 'level': 1},
    {'text': 'Upside Case (Year 4): 50-60x return | $3-4B exit', 'size': 14, 'level': 1},
    {'text': 'IPO Case (Year 6-7): 100-150x return | $5-10B valuation', 'size': 14, 'level': 1},
])

# SLIDE 14: Risk Mitigation
print("  → Slide 14: Risk Mitigation")
add_content_slide(prs, "Risk Mitigation", [
    {'text': 'Market Adoption Risk: Early customer validation, freemium model, partnerships', 'size': 16},
    {'text': 'Churn & Retention: Net revenue retention >110%, dedicated CSM', 'size': 16},
    {'text': 'Technology Risk: 95% accuracy proven, redundant architecture', 'size': 16},
    {'text': 'Sales Execution: Experienced VP Sales hire, founder sales track record', 'size': 16},
    {'text': 'Funding Risk: 18+ months runway on $2M, profitable by Year 1', 'size': 16},
    {'text': 'Economic Downturn: ROI-positive, compliance-driven spend continues', 'size': 16},
])

# SLIDE 15: The Pitch (3 minutes)
print("  → Slide 15: The 3-Minute Pitch")
add_content_slide(prs, "The Investment Thesis", [
    {'text': '$2.5T market', 'size': 20, 'bold': True, 'color': COLOR_ACCENT},
    {'text': 'Zero competition | Structural moat | Exceptional unit economics', 'size': 16, 'level': 1},
    {'text': '', 'size': 12},
    {'text': 'Solving supply chain intelligence equity', 'size': 20, 'bold': True, 'color': COLOR_ACCENT},
    {'text': 'Manifesto-driven mission meets venture opportunity', 'size': 16, 'level': 1},
    {'text': '', 'size': 12},
    {'text': 'Profitable by Year 1', 'size': 20, 'bold': True, 'color': COLOR_ACCENT},
    {'text': '55% EBITDA margin Year 1 | 90% gross margin | Capital efficient', 'size': 16, 'level': 1},
])

# SLIDE 16: Call to Action
print("  → Slide 16: Call to Action")
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARY

cta_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2.5))
cta_frame = cta_box.text_frame
cta_frame.word_wrap = True

p = cta_frame.paragraphs[0]
p.text = "Category Creation Opportunity"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.alignment = PP_ALIGN.CENTER

cta_frame.add_paragraph()
p = cta_frame.paragraphs[1]
p.text = "Let's build intelligence equity together"
p.font.size = Pt(32)
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

cta_frame.add_paragraph()
p = cta_frame.paragraphs[2]
p.text = ""

cta_frame.add_paragraph()
p = cta_frame.paragraphs[3]
p.text = "Contact: contact@kraftd.io | Demo: demo.kraftd.io"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r'c:\Users\1R6\OneDrive\Project Catalyst\KraftdIntel\KRAFTD_Strategic_Pitch_Deck.pptx'
prs.save(output_path)
print(f'\n✓ PowerPoint presentation created: {output_path}')
print(f'✓ 16 professional slides with investor-grade design')
print(f'✓ Brand colors: Dark Blue (#003366) + Orange (#FF6600)')
print(f'✓ Includes: Problem, Market, Solution, Business Model, Financials, GTM, Risks, CTA')
print(f'✓ Ready for investor presentations')
